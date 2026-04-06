#!/usr/bin/env python3
"""
TC Management System Presentation Builder
발표용 슬라이드 생성 스크립트 (python-pptx 기반)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os
import datetime

# ── Color Palette ──────────────────────────────────────
C_NAVY   = RGBColor(0x1B, 0x2E, 0x4B)   # dark navy
C_BLUE   = RGBColor(0x23, 0x6F, 0xAB)   # mid blue
C_TEAL   = RGBColor(0x00, 0xA8, 0xA8)   # teal accent
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)   # green
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)   # orange
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)   # red
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY   = RGBColor(0x74, 0x7D, 0x8C)
C_LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)
C_DIVIDER= RGBColor(0xD5, 0xDB, 0xE5)
C_YELLOW = RGBColor(0xF3, 0x9C, 0x12)

# ── Dimensions ────────────────────────────────────────
W = Inches(10)   # slide width
H = Inches(7.5)  # slide height
MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.5)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=18, font_color=C_BLACK,
                bold=False, italic=False, align=PP_ALIGN.LEFT,
                font_name="Apple SD Gothic Neo", wrap=True, v_anchor=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox


def add_text_frame(slide, x, y, w, h, lines, default_size=16,
                   default_color=C_BLACK, font_name="Apple SD Gothic Neo",
                   wrap=True, line_spacing=1.2):
    """
    lines: list of (text, size, color, bold, align, indent_level)
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap

    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, color, bold, align, level = item, default_size, default_color, False, PP_ALIGN.LEFT, 0
        elif len(item) == 2:
            text, size = item; color, bold, align, level = default_color, False, PP_ALIGN.LEFT, 0
        elif len(item) == 3:
            text, size, color = item; bold, align, level = False, PP_ALIGN.LEFT, 0
        elif len(item) == 4:
            text, size, color, bold = item; align, level = PP_ALIGN.LEFT, 0
        elif len(item) == 5:
            text, size, color, bold, align = item; level = 0
        else:
            text, size, color, bold, align, level = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        p.level = level

        # line spacing
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsmap
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))

        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name

    return txBox


def add_title_bar(slide, title_text, bg_color=C_NAVY, text_color=C_WHITE,
                  bar_height=Inches(1.1), subtitle=None):
    """Add standard title bar at top of slide."""
    # background bar
    add_rect(slide, 0, 0, W, bar_height, fill_color=bg_color)
    # accent line at bottom of bar
    add_rect(slide, 0, bar_height - Inches(0.04), W, Inches(0.04), fill_color=C_TEAL)

    # title text
    txBox = slide.shapes.add_textbox(MARGIN, Inches(0.18), W - MARGIN * 2, Inches(0.75))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Apple SD Gothic Neo"

    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.color.rgb = C_TEAL
        r2.font.name = "Apple SD Gothic Neo"


def bullet_block(slide, x, y, w, h, bullets, icon_color=C_BLUE,
                 title=None, title_color=C_NAVY, bg_color=None, padding=Inches(0.15)):
    """Draw a content block with optional title and bullet list."""
    if bg_color:
        add_rect(slide, x, y, w, h, fill_color=bg_color)

    cy = y + padding
    if title:
        tb = slide.shapes.add_textbox(x + padding, cy, w - padding*2, Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Apple SD Gothic Neo"
        cy += Inches(0.45)

    tb = slide.shapes.add_textbox(x + padding, cy, w - padding*2, h - (cy - y) - padding)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if isinstance(b, str):
            text, size, color, bold = b, 15, C_BLACK, False
        elif len(b) == 2:
            text, size = b; color, bold = C_BLACK, False
        elif len(b) == 3:
            text, size, color = b; bold = False
        else:
            text, size, color, bold = b

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Apple SD Gothic Neo"


def step_box(slide, x, y, w, h, number, title, desc, bg_color=C_BLUE,
             role_tag=None, role_color=None):
    """Draw a numbered step box."""
    add_rect(slide, x, y, w, h, fill_color=bg_color)

    # circle number background
    add_rect(slide, x + Inches(0.12), y + Inches(0.1), Inches(0.38), Inches(0.38),
             fill_color=C_WHITE)

    # number
    add_textbox(slide, x + Inches(0.12), y + Inches(0.08), Inches(0.38), Inches(0.38),
                str(number), font_size=18, font_color=bg_color, bold=True,
                align=PP_ALIGN.CENTER)

    # role tag
    tx = x + Inches(0.6)
    ty = y + Inches(0.1)
    if role_tag:
        tb = slide.shapes.add_textbox(tx, ty, w - Inches(0.7), Inches(0.28))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = role_tag
        run.font.size = Pt(11)
        run.font.color.rgb = role_color or C_WHITE
        run.font.name = "Apple SD Gothic Neo"
        run.font.bold = True
        ty += Inches(0.28)

    # title
    tb = slide.shapes.add_textbox(tx, ty, w - Inches(0.7), Inches(0.35))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.color.rgb = C_WHITE
    run.font.bold = True
    run.font.name = "Apple SD Gothic Neo"

    if desc:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = RGBColor(0xD0, 0xE8, 0xFF)
        r2.font.name = "Apple SD Gothic Neo"


def arrow(slide, x, y, vertical=True, color=C_GRAY):
    """Draw a small arrow."""
    size = Inches(0.25)
    if vertical:
        add_textbox(slide, x - size/2, y, size, size, "▼",
                    font_size=14, font_color=color, align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, x, y - size/2, size, size, "▶",
                    font_size=14, font_color=color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Cover: 테스트 케이스 관리 시스템 제안"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # full background
    add_rect(slide, 0, 0, W, H, fill_color=C_NAVY)
    # bottom accent stripe
    add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), fill_color=C_TEAL)
    # left accent bar
    add_rect(slide, 0, 0, Inches(0.35), H, fill_color=C_BLUE)

    # tag line
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(8), Inches(0.5),
                "QA TEAM  ·  TC MANAGEMENT SYSTEM PROPOSAL",
                font_size=12, font_color=C_TEAL, bold=True)

    # main title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(8.5), Inches(1.8))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "테스트 케이스 관리 시스템 제안"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = C_WHITE
    run.font.name = "Apple SD Gothic Neo"

    # subtitle
    add_textbox(slide, Inches(0.7), Inches(3.95), Inches(8), Inches(0.55),
                "TC 작성 · 관리 · 실행 · 결과까지 — 단일 통합 시스템 설계",
                font_size=18, font_color=RGBColor(0xA8, 0xD8, 0xEA))

    # divider
    add_rect(slide, Inches(0.7), Inches(4.55), Inches(4), Inches(0.04),
             fill_color=C_TEAL)

    # version / date
    today = datetime.date.today().strftime("%Y.%m.%d")
    add_textbox(slide, Inches(0.7), Inches(4.75), Inches(8), Inches(0.4),
                f"v0.2   |   {today}",
                font_size=13, font_color=C_GRAY)

    # decorative dots (top-right)
    for i in range(5):
        for j in range(5):
            alpha = 0x20 + (i + j) * 0x08
            c = RGBColor(min(alpha, 0xFF), min(alpha + 0x30, 0xFF), min(alpha + 0x60, 0xFF))
            add_rect(slide,
                     Inches(8.0 + j * 0.25), Inches(0.5 + i * 0.25),
                     Inches(0.18), Inches(0.18), fill_color=c)


def slide_02_current_state(prs):
    """현황 & 과제"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "TC 자동화 개발 — 현황 & 과제")

    # ── LEFT: 현황 ──
    bx, by, bw, bh = MARGIN, CONTENT_TOP, Inches(4.1), Inches(5.0)
    add_rect(slide, bx, by, bw, bh, fill_color=C_WHITE)
    add_rect(slide, bx, by, bw, Inches(0.45), fill_color=C_BLUE)
    add_textbox(slide, bx + Inches(0.2), by + Inches(0.05), bw - Inches(0.3), Inches(0.4),
                "📋  현황", font_size=16, font_color=C_WHITE, bold=True)

    bullets = [
        ("전체 TC: v2.0 기준으로 정리 완료", 15, C_BLACK, False),
        ("  → v3.0에서 자동화 여부 기록 완료", 13, C_GRAY, False),
        ("", 8, C_BLACK, False),
        ("Phase 구분", 15, C_BLUE, True),
        ("  P1: Youthmeta 거래소 (5개 거래소)", 14, C_BLACK, False),
        ("  P2: Mobile 앱", 14, C_BLACK, False),
        ("", 8, C_BLACK, False),
        ("TC 형식: BDD 기반, Excel/Markdown 병행 관리", 14, C_BLACK, False),
    ]
    bullet_block(slide, bx, by + Inches(0.5), bw, bh - Inches(0.5), bullets,
                 padding=Inches(0.2))

    # ── RIGHT: 남은 과제 ──
    rx, ry, rw, rh = Inches(5.0), CONTENT_TOP, Inches(4.4), Inches(5.0)
    add_rect(slide, rx, ry, rw, rh, fill_color=C_WHITE)
    add_rect(slide, rx, ry, rw, Inches(0.45), fill_color=C_ORANGE)
    add_textbox(slide, rx + Inches(0.2), ry + Inches(0.05), rw - Inches(0.3), Inches(0.4),
                "⚠️  필요한 작업", font_size=16, font_color=C_WHITE, bold=True)

    tasks = [
        ("① TC ID 체계 정비", 15, C_NAVY, True),
        ("    SPCY-{Phase}-{번호} 형식 전체 소급 적용", 13, C_GRAY, False),
        ("", 8, C_BLACK, False),
        ("② TC 추가 작성", 15, C_NAVY, True),
        ("    Youthmeta / Mobile TC 신규 작성 필요", 13, C_GRAY, False),
        ("", 8, C_BLACK, False),
        ("③ Phase 1 TC 리뷰", 15, C_NAVY, True),
        ("    거래소별 TC 커버리지 검증 필요", 13, C_GRAY, False),
        ("", 8, C_BLACK, False),
        ("④ 자동화 연동", 15, C_NAVY, True),
        ("    is_automated / auto_tc_id 필드 추가", 13, C_GRAY, False),
    ]
    bullet_block(slide, rx, ry + Inches(0.5), rw, rh - Inches(0.5), tasks,
                 padding=Inches(0.2))


def slide_03_before_after(prs):
    """Before vs After"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "Before vs After — 시스템 도입 효과")

    col_w = Inches(4.3)
    gap = Inches(0.2)
    left_x = MARGIN
    right_x = MARGIN + col_w + gap
    top_y = CONTENT_TOP
    col_h = Inches(5.2)

    # BEFORE column
    add_rect(slide, left_x, top_y, col_w, col_h, fill_color=C_WHITE)
    add_rect(slide, left_x, top_y, col_w, Inches(0.48), fill_color=RGBColor(0xB7, 0x3A, 0x3A))
    add_textbox(slide, left_x + Inches(0.2), top_y + Inches(0.06),
                col_w - Inches(0.3), Inches(0.38),
                "✗  Before", font_size=17, font_color=C_WHITE, bold=True)

    before_items = [
        "TC 작성: 수동 (QA 개인 역량 의존)",
        "TC 형식: 작성자마다 다름",
        "TC 저장: 개별 Excel, 네이밍 비일관",
        "TC ID: 없거나 파일마다 달라 추적 불가",
        "자동화 연동: TC와 자동화 코드 분리",
        "결과 기록: 실행 후 별도 업데이트",
        "현황 파악: 파일 열어봐야 알 수 있음",
        "리포트: 별도 수작업 작성",
    ]
    cy = top_y + Inches(0.55)
    for item in before_items:
        add_textbox(slide, left_x + Inches(0.25), cy, col_w - Inches(0.4), Inches(0.42),
                    f"  {item}", font_size=13, font_color=RGBColor(0x55, 0x55, 0x55))
        cy += Inches(0.52)

    # AFTER column
    add_rect(slide, right_x, top_y, col_w, col_h, fill_color=C_WHITE)
    add_rect(slide, right_x, top_y, col_w, Inches(0.48), fill_color=C_GREEN)
    add_textbox(slide, right_x + Inches(0.2), top_y + Inches(0.06),
                col_w - Inches(0.3), Inches(0.38),
                "✓  After", font_size=17, font_color=C_WHITE, bold=True)

    after_items = [
        "Claude + Harness 자동 초안 생성",
        "BDD / 규칙 파일 기준 통일",
        "SPCY_TC_{Phase}_{날짜}_v{N}.xlsx",
        "SPCY-{Phase}-{번호} 고유 식별자",
        "is_automated + auto_tc_id 연동",
        "수동 기록 + 자동 TC 결과 자동 연동",
        "대시보드에서 실시간 진행률 확인",
        "Test Result Agent 기반 자동 생성",
    ]
    cy = top_y + Inches(0.55)
    for item in after_items:
        add_textbox(slide, right_x + Inches(0.25), cy, col_w - Inches(0.4), Inches(0.42),
                    f"  {item}", font_size=13, font_color=RGBColor(0x22, 0x55, 0x33))
        cy += Inches(0.52)

    # center VS badge
    vs_x = MARGIN + col_w + gap / 2 - Inches(0.22)
    vs_y = top_y + col_h / 2 - Inches(0.3)
    add_rect(slide, vs_x, vs_y, Inches(0.44), Inches(0.6), fill_color=C_NAVY)
    add_textbox(slide, vs_x, vs_y + Inches(0.05), Inches(0.44), Inches(0.5),
                "VS", font_size=14, font_color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


def slide_04_writing_process(prs):
    """TC 작성 프로세스"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "TC 작성 프로세스")

    steps = [
        (C_BLUE,   "초기 작성",  "TC Generator 실행\n(Claude + Harness 파이프라인)"),
        (C_TEAL,   "QA 리뷰",   "QA 팀 검토\n도메인별 커버리지 확인"),
        (C_GREEN,  "수정 반영",  "Generator 기반 자동 재작성\n(누락 기능 보강)"),
        (C_ORANGE, "산출물 배포", "Excel / Markdown 생성\nSPCY_TC_{Phase}_{날짜}_v{N}"),
        (C_NAVY,   "버전 관리",  "Phase / Sprint 단위 관리\nMinor 버전 증가"),
    ]

    sx = MARGIN
    step_w = Inches(1.55)
    step_h = Inches(2.8)
    sy = Inches(1.9)
    gap = Inches(0.2)

    for i, (color, title, desc) in enumerate(steps):
        x = sx + i * (step_w + gap)
        add_rect(slide, x, sy, step_w, step_h, fill_color=color)

        # number circle
        add_rect(slide, x + Inches(0.55), sy + Inches(0.18), Inches(0.45), Inches(0.45),
                 fill_color=C_WHITE)
        add_textbox(slide, x + Inches(0.55), sy + Inches(0.15), Inches(0.45), Inches(0.45),
                    str(i+1), font_size=20, font_color=color, bold=True,
                    align=PP_ALIGN.CENTER)

        # title
        add_textbox(slide, x + Inches(0.1), sy + Inches(0.75), step_w - Inches(0.15),
                    Inches(0.5), title, font_size=15, font_color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

        # desc
        tb = slide.shapes.add_textbox(x + Inches(0.1), sy + Inches(1.3),
                                      step_w - Inches(0.15), Inches(1.3))
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = desc
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor(0xD8, 0xEC, 0xFF)
        r.font.name = "Apple SD Gothic Neo"

        # arrow
        if i < len(steps) - 1:
            ax = x + step_w + gap / 2 - Inches(0.12)
            add_textbox(slide, ax, sy + step_h / 2 - Inches(0.2), Inches(0.25), Inches(0.4),
                        "▶", font_size=16, font_color=C_GRAY, align=PP_ALIGN.CENTER)

    # bottom note
    add_textbox(slide, MARGIN, Inches(5.05), W - MARGIN*2, Inches(0.5),
                "결과물 형식: Excel (.xlsx) + Markdown (.md)  |  파일명: SPCY_TC_{Phase}_{YYYYMMDD}_v{N}",
                font_size=12, font_color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_05_management_structure(prs):
    """TC 관리 구조"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "TC 관리 구조")

    # ── Top box: Master TC ──
    mx, my, mw, mh = Inches(3.0), CONTENT_TOP, Inches(4.0), Inches(1.0)
    add_rect(slide, mx, my, mw, mh, fill_color=C_NAVY)
    add_textbox(slide, mx, my + Inches(0.1), mw, Inches(0.38),
                "📦  Master TC Repository", font_size=16, font_color=C_WHITE, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, mx, my + Inches(0.48), mw, Inches(0.38),
                "전체 TC 통합 관리  |  TC ID: UNIQUE  |  버전 관리",
                font_size=12, font_color=C_TEAL, align=PP_ALIGN.CENTER)

    # down arrow
    add_textbox(slide, Inches(4.8), my + mh + Inches(0.02), Inches(0.4), Inches(0.35),
                "▼", font_size=16, font_color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── Phase boxes ──
    phases = [
        (C_BLUE,   "Phase 1\nYouthmeta", "5개 거래소\nP1-{번호}"),
        (C_TEAL,   "Phase 2\nMobile",    "iOS / Android\nP2-{번호}"),
        (C_GREEN,  "Phase N\n신규 추가",  "확장 가능 구조\nPN-{번호}"),
    ]
    pw, ph = Inches(2.6), Inches(1.4)
    pgap = Inches(0.35)
    total_pw = len(phases) * pw + (len(phases)-1) * pgap
    px_start = (W - total_pw) / 2
    py = Inches(3.3)

    for i, (color, title, detail) in enumerate(phases):
        px = px_start + i * (pw + pgap)
        add_rect(slide, px, py, pw, ph, fill_color=color)
        add_textbox(slide, px + Inches(0.1), py + Inches(0.1), pw - Inches(0.2),
                    Inches(0.55), title, font_size=14, font_color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, px + Inches(0.1), py + Inches(0.7), pw - Inches(0.2),
                    Inches(0.6), detail, font_size=12,
                    font_color=RGBColor(0xD0, 0xF0, 0xFF), align=PP_ALIGN.CENTER)

    # ── Bottom: Key principles ──
    by = Inches(5.05)
    items = [
        ("TC ID는 고유 식별자", C_NAVY),
        ("Phase별 독립 관리", C_BLUE),
        ("자동화 여부 포함", C_TEAL),
        ("버전 히스토리 보존", C_GREEN),
    ]
    iw = Inches(2.0)
    igap = Inches(0.25)
    total_iw = len(items) * iw + (len(items)-1) * igap
    ix_start = (W - total_iw) / 2
    for i, (text, color) in enumerate(items):
        ix = ix_start + i * (iw + igap)
        add_rect(slide, ix, by, iw, Inches(0.5), fill_color=color)
        add_textbox(slide, ix, by + Inches(0.05), iw, Inches(0.4),
                    text, font_size=12, font_color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)


def slide_06_automation_process(prs):
    """TC 자동화 개발 프로세스 — 5단계"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "TC 자동화 개발 프로세스")

    steps = [
        (C_BLUE,   "[QA]",      "Phase별 TC 초안 작성",         "Claude + Harness 자동 생성"),
        (C_TEAL,   "[QA → 개발]", "TC 전달 + 자동화 여부 공유",  "TC ID 기반으로 자동화 대상 명시"),
        (C_GREEN,  "[개발]",    "자동화 케이스 확정",            "구현 가능 여부 판단"),
        (C_ORANGE, "[QA]",      "자동 / 수동 분류 및 스케줄 확정", "실행 계획 수립"),
        (C_NAVY,   "[QA + 개발]", "테스트 실행 ↔ 자동화 구현",    "병렬 진행, 결과 통합"),
    ]

    step_h = Inches(0.85)
    sy = CONTENT_TOP
    sx = MARGIN
    sw = W - MARGIN * 2

    for i, (color, role, title, desc) in enumerate(steps):
        y = sy + i * (step_h + Inches(0.08))
        add_rect(slide, sx, y, sw, step_h, fill_color=color)

        # step number
        add_rect(slide, sx, y, Inches(0.55), step_h, fill_color=RGBColor(0, 0, 0))
        # make slightly transparent by using a darker shade
        add_textbox(slide, sx, y + Inches(0.18), Inches(0.55), Inches(0.5),
                    str(i+1), font_size=22, font_color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

        # role tag
        add_textbox(slide, sx + Inches(0.65), y + Inches(0.08), Inches(1.4), Inches(0.32),
                    role, font_size=11, font_color=RGBColor(0xFF, 0xFF, 0xCC), bold=True)

        # title
        add_textbox(slide, sx + Inches(0.65), y + Inches(0.38), Inches(4.0), Inches(0.42),
                    title, font_size=15, font_color=C_WHITE, bold=True)

        # desc
        add_textbox(slide, sx + Inches(4.8), y + Inches(0.2), Inches(4.5), Inches(0.5),
                    desc, font_size=13, font_color=RGBColor(0xD8, 0xEF, 0xFF))

        # arrow (except last)
        if i < len(steps) - 1:
            ay = y + step_h + Inches(0.01)
            add_textbox(slide, sx + sw/2 - Inches(0.2), ay, Inches(0.4), Inches(0.1),
                        "▼", font_size=10, font_color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_07_execution(prs):
    """TC 실행 & 결과 관리"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "TC 실행 & 결과 관리")

    col_w = Inches(4.3)
    gap = Inches(0.3)
    top_y = CONTENT_TOP
    col_h = Inches(5.1)

    # LEFT: 실행 프로세스
    lx = MARGIN
    add_rect(slide, lx, top_y, col_w, col_h, fill_color=C_WHITE)
    add_rect(slide, lx, top_y, col_w, Inches(0.45), fill_color=C_BLUE)
    add_textbox(slide, lx + Inches(0.2), top_y + Inches(0.05),
                col_w - Inches(0.3), Inches(0.38),
                "▶  TC 실행 프로세스", font_size=15, font_color=C_WHITE, bold=True)

    exec_items = [
        ("수동 TC", "QA에서 버전별 Excel에 결과 직접 기록", C_BLUE),
        ("자동 TC", "자동화 시스템이 결과 자동 기록\n(is_automated = true)", C_TEAL),
        ("실행 중 변경", "TC 추가 / 수정 / 삭제 가능\n(Minor 버전 올려 관리)", C_ORANGE),
    ]
    cy = top_y + Inches(0.6)
    for label, detail, color in exec_items:
        add_rect(slide, lx + Inches(0.2), cy, Inches(0.06), Inches(0.85), fill_color=color)
        add_textbox(slide, lx + Inches(0.4), cy, col_w - Inches(0.55), Inches(0.32),
                    label, font_size=14, font_color=color, bold=True)
        tb = slide.shapes.add_textbox(lx + Inches(0.4), cy + Inches(0.32),
                                      col_w - Inches(0.55), Inches(0.55))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = detail; r.font.size = Pt(12)
        r.font.color.rgb = C_GRAY; r.font.name = "Apple SD Gothic Neo"
        cy += Inches(1.1)

    # RIGHT: 결과 관리
    rx = MARGIN + col_w + gap
    add_rect(slide, rx, top_y, col_w, col_h, fill_color=C_WHITE)
    add_rect(slide, rx, top_y, col_w, Inches(0.45), fill_color=C_TEAL)
    add_textbox(slide, rx + Inches(0.2), top_y + Inches(0.05),
                col_w - Inches(0.3), Inches(0.38),
                "📊  결과 관리", font_size=15, font_color=C_WHITE, bold=True)

    result_items = [
        ("목표 일정 관리", "시험 기간 설정 및 진행률 추적", C_NAVY),
        ("실시간 모니터링", "Pass / Fail / N/T 현황 대시보드", C_BLUE),
        ("자동화팀 역할", "이슈 디버깅 & 케이스 보고", C_TEAL),
        ("결과서 자동 생성", "Test Result Agent 기반\nSlack / 메일 자동 발송", C_GREEN),
    ]
    cy = top_y + Inches(0.6)
    for label, detail, color in result_items:
        add_rect(slide, rx + Inches(0.2), cy, Inches(0.06), Inches(0.85), fill_color=color)
        add_textbox(slide, rx + Inches(0.4), cy, col_w - Inches(0.55), Inches(0.32),
                    label, font_size=14, font_color=color, bold=True)
        tb = slide.shapes.add_textbox(rx + Inches(0.4), cy + Inches(0.32),
                                      col_w - Inches(0.55), Inches(0.55))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = detail; r.font.size = Pt(12)
        r.font.color.rgb = C_GRAY; r.font.name = "Apple SD Gothic Neo"
        cy += Inches(1.1)


def slide_08_dashboard(prs):
    """대시보드 기능"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "대시보드 기능")

    cards = [
        (C_RED,    "🔴  Fail 리스트",
         "• 중요도별 목록 (High → Low)\n• 자동 리포트 → 메일 / 슬랙"),
        (C_BLUE,   "🔍  검색 & 필터",
         "• 거래소별 결과 필터링\n• 키워드 전체 검색"),
        (C_GREEN,  "📈  통계 & 진행률",
         "• Pass/Fail/N/T 실시간 현황\n• 버전별 비교"),
        (C_ORANGE, "🤖  자동화 비율",
         "• 전체 대비 자동화 TC 비율\n• Phase별 자동화 현황"),
    ]

    cw = Inches(4.3)
    ch = Inches(2.2)
    gap = Inches(0.3)
    for i, (color, title, detail) in enumerate(cards):
        row, col = divmod(i, 2)
        cx = MARGIN + col * (cw + gap)
        cy = CONTENT_TOP + row * (ch + gap)
        add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        # left color bar
        add_rect(slide, cx, cy, Inches(0.1), ch, fill_color=color)
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.3), Inches(0.45),
                    title, font_size=15, font_color=color, bold=True)
        tb = slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.65),
                                      cw - Inches(0.3), Inches(1.4))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = detail; r.font.size = Pt(13)
        r.font.color.rgb = C_GRAY; r.font.name = "Apple SD Gothic Neo"


def slide_09_requirements(prs):
    """시스템 요구사항"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "시스템 요구사항")

    reqs = [
        (C_BLUE,   "👁  Viewer 기능",
         "• 컬럼별 선택 옵션으로 TC 보기\n• Excel / CSV export\n• 거래소 / Phase / 도메인 필터"),
        (C_TEAL,   "✏️  Edit 기능",
         "• Edit → 임시저장 → 승인 → 확정 후 TC 업데이트\n• 승인 단계 Skip 옵션\n• 변경 히스토리 자동 기록"),
        (C_GREEN,  "👥  협업 환경",
         "• 1인 작업 또는 팀 작업 모드\n• 동시 편집 지원\n• 권한 관리"),
        (C_ORANGE, "🔗  자동화 연동",
         "• 자동화 시스템과 TC ID 연동\n• is_automated / auto_tc_id 필드\n• 결과 자동 수집 API"),
    ]

    cw = Inches(4.3)
    ch = Inches(2.2)
    gap = Inches(0.3)
    for i, (color, title, detail) in enumerate(reqs):
        row, col = divmod(i, 2)
        cx = MARGIN + col * (cw + gap)
        cy = CONTENT_TOP + row * (ch + gap)
        add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        add_rect(slide, cx, cy, Inches(0.1), ch, fill_color=color)
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.3), Inches(0.45),
                    title, font_size=15, font_color=color, bold=True)
        tb = slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.65),
                                      cw - Inches(0.3), Inches(1.4))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = detail; r.font.size = Pt(13)
        r.font.color.rgb = C_GRAY; r.font.name = "Apple SD Gothic Neo"


def slide_10_tech_stack(prs):
    """기술 스택"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "기술 스택")

    stacks = [
        (C_NAVY,   "TC 생성",
         ["Claude Sonnet (AI 모델)",
          "Harness (에이전트 오케스트레이터)",
          "파이프라인: doc-parser → tc-writer → tc-reviewer",
          "규칙: common/tc-rules.md + tc-rules-override.md"]),
        (C_BLUE,   "TC 저장 / 배포",
         ["Excel (.xlsx) + Markdown (.md) 병행",
          "파일명: SPCY_TC_{Phase}_{날짜}_v{N}",
          "빌더: scripts/build_excel.py (openpyxl)",
          "버전 자동 증가 관리"]),
        (C_TEAL,   "TC 조회 / 실행",
         ["Viewer: 웹 기반 (구현 예정)",
          "컬럼 선택 / 거래소 필터 / Export",
          "수동 실행: Excel 직접 결과 기록",
          "자동 실행: is_automated / auto_tc_id"]),
        (C_GREEN,  "결과 리포팅",
         ["Test Result Agent (자동 결과서 생성)",
          "알림: Slack / 메일 자동 리포트",
          "Fail 중요도별 분류",
          "진행률 실시간 대시보드"]),
    ]

    sw = Inches(2.0)
    sh = Inches(4.8)
    gap = Inches(0.25)
    sx = MARGIN + Inches(0.1)

    for i, (color, title, items) in enumerate(stacks):
        x = sx + i * (sw + gap)
        y = CONTENT_TOP

        add_rect(slide, x, y, sw, sh, fill_color=C_WHITE)
        add_rect(slide, x, y, sw, Inches(0.5), fill_color=color)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.06), sw - Inches(0.15),
                    Inches(0.38), title, font_size=14, font_color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

        cy = y + Inches(0.65)
        for item in items:
            tb = slide.shapes.add_textbox(x + Inches(0.12), cy, sw - Inches(0.2), Inches(0.85))
            tb.word_wrap = True
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = f"• {item}"; r.font.size = Pt(11)
            r.font.color.rgb = C_BLACK; r.font.name = "Apple SD Gothic Neo"
            cy += Inches(0.92)


def slide_11_data_structure(prs):
    """TC 데이터 구조 설계"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "TC 데이터 구조 설계")

    groups = [
        (C_NAVY,   "기본 식별",
         "tc_id  |  phase  |  domain  |  version"),
        (C_BLUE,   "테스트 내용",
         "title  /  priority (H/M/L)  /  type (Positive/Negative/Edge)\nprecondition  /  given · when · then  /  expected_result  /  platform"),
        (C_TEAL,   "거래소 필드  (Phase 1)",
         "Gate · OKX · Bybit · Bitget · Hyperliquid\n→ Pass / Fail / N/T / N/A  |  tester: Tester 1 / Tester 2"),
        (C_ORANGE, "자동화 연동",
         "is_automated: true/false\nauto_tc_id: 자동화 시스템 연동 ID\nautomation_status: 완료 / 진행중 / 미정"),
        (C_GREEN,  "결과 & 이력",
         "status: Pass/Fail/N/T/N/A\nactual_result  /  tested_at  /  tester\nchange_log: [ { version, who, what, when } ]"),
    ]

    gw = W - MARGIN * 2
    gh = Inches(0.82)
    gap = Inches(0.1)
    gy = CONTENT_TOP

    for color, title, fields in groups:
        add_rect(slide, MARGIN, gy, Inches(0.08), gh, fill_color=color)
        add_rect(slide, MARGIN + Inches(0.12), gy, gw - Inches(0.12), gh, fill_color=C_WHITE)

        add_textbox(slide, MARGIN + Inches(0.25), gy + Inches(0.06),
                    Inches(1.8), Inches(0.32), title,
                    font_size=13, font_color=color, bold=True)

        tb = slide.shapes.add_textbox(MARGIN + Inches(2.1), gy + Inches(0.06),
                                      gw - Inches(2.3), gh - Inches(0.12))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = fields; r.font.size = Pt(11)
        r.font.color.rgb = C_GRAY; r.font.name = "Courier New"
        gy += gh + gap


def slide_12_change_management(prs):
    """변경 관리"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "변경 관리")

    phases_data = [
        (C_BLUE,   "테스트 시작 전",
         ["TC 수정 시 자동화팀 공유",
          "자동화 여부 컨펌 프로세스 적용",
          "전체 TC에 반영 (수시 또는 일괄)"]),
        (C_ORANGE, "테스트 진행 중",
         ["기본 원칙: Minor 버전 올린 후 tc id 관리",
          "스텝 / 기대결과 변경 내역 기록",
          "긴급 수정: 즉시 공유 → 반영 → 재진행"]),
        (C_GREEN,  "변경 히스토리 (Optional)",
         ["tc id + 수정 내용 기록",
          "변경 날짜 & 작성자 보존",
          "change_log 필드로 관리"]),
    ]

    cw = Inches(2.75)
    ch = Inches(4.4)
    gap = Inches(0.35)
    total_w = len(phases_data) * cw + (len(phases_data)-1) * gap
    cx_start = (W - total_w) / 2
    cy = CONTENT_TOP

    for i, (color, title, items) in enumerate(phases_data):
        cx = cx_start + i * (cw + gap)
        add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        add_rect(slide, cx, cy, cw, Inches(0.5), fill_color=color)
        add_textbox(slide, cx + Inches(0.15), cy + Inches(0.06), cw - Inches(0.25),
                    Inches(0.38), title, font_size=14, font_color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

        item_y = cy + Inches(0.65)
        for item in items:
            tb = slide.shapes.add_textbox(cx + Inches(0.2), item_y,
                                          cw - Inches(0.35), Inches(0.9))
            tb.word_wrap = True
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = f"• {item}"; r.font.size = Pt(13)
            r.font.color.rgb = C_BLACK; r.font.name = "Apple SD Gothic Neo"
            item_y += Inches(1.0)


def slide_13_conclusion(prs):
    """결론"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_NAVY)
    add_rect(slide, 0, H - Inches(0.5), W, Inches(0.5), fill_color=C_TEAL)
    add_rect(slide, 0, 0, Inches(0.35), H, fill_color=C_BLUE)

    add_textbox(slide, Inches(0.7), Inches(0.9), Inches(8), Inches(0.5),
                "CONCLUSION", font_size=13, font_color=C_TEAL, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(8), Inches(0.7),
                "결론", font_size=36, font_color=C_WHITE, bold=True)
    add_rect(slide, Inches(0.7), Inches(2.05), Inches(3.5), Inches(0.04),
             fill_color=C_TEAL)

    points = [
        ("AI 기반 TC 자동 생성",   "Claude + Harness로 TC 초안을 자동 생성, QA 리소스 집중 가능",   C_BLUE),
        ("통합 관리 & 추적",       "Master TC + Phase별 관리로 누락·중복 없이 체계적 추적",         C_TEAL),
        ("자동화 연동",            "개발팀과 TC ID 기반으로 협업, 자동화 결과 자동 수집",            C_GREEN),
        ("실시간 대시보드",         "진행률·Fail 현황을 즉시 파악, 리포트 자동 발송",               C_ORANGE),
        ("QA 생산성 극대화",        "수동 작업 최소화, 단일 시스템으로 전 과정 통합",               C_YELLOW),
    ]

    py = Inches(2.3)
    for label, desc, color in points:
        add_rect(slide, Inches(0.7), py, Inches(0.08), Inches(0.5), fill_color=color)
        add_textbox(slide, Inches(0.9), py, Inches(2.5), Inches(0.32),
                    label, font_size=14, font_color=color, bold=True)
        add_textbox(slide, Inches(0.9), py + Inches(0.3), Inches(7.8), Inches(0.32),
                    desc, font_size=13, font_color=RGBColor(0xCC, 0xD6, 0xE8))
        py += Inches(0.85)


def slide_14_next_steps(prs):
    """Next Steps (로드맵)"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill_color=C_LGRAY)
    add_title_bar(slide, "Next Steps  —  로드맵")

    phases = [
        (C_RED,    "즉시",   "TC 데이터 구조 확정",
         ["SPCY-{Phase}-{번호} 포맷 전체 소급",
          "Youthmeta / Mobile TC 마이그레이션"]),
        (C_ORANGE, "단기",   "TC 관리 시스템 MVP",
         ["Viewer: 컬럼 선택 + 거래소 필터 + Export",
          "Edit 기능: 임시저장 → 승인 플로우",
          "자동화 TC ID 연동 (is_automated)"]),
        (C_BLUE,   "중기",   "실행 & 결과 관리 연동",
         ["수동 / 자동 TC 결과 기록 통합",
          "Fail 리스트 → 슬랙/메일 자동 리포트"]),
        (C_GREEN,  "장기",   "대시보드 & Test Result Agent",
         ["진행률 / 자동화 비율 / 거래소별 현황",
          "Test Result Agent 기반 결과서 자동 생성"]),
    ]

    cw = Inches(2.0)
    ch = Inches(4.6)
    gap = Inches(0.3)
    total_w = len(phases) * cw + (len(phases)-1) * gap
    cx_start = (W - total_w) / 2
    cy = CONTENT_TOP

    for i, (color, timing, title, items) in enumerate(phases):
        cx = cx_start + i * (cw + gap)
        add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        add_rect(slide, cx, cy, cw, Inches(0.75), fill_color=color)

        add_textbox(slide, cx + Inches(0.1), cy + Inches(0.04), cw - Inches(0.15),
                    Inches(0.3), timing, font_size=12,
                    font_color=RGBColor(0xFF, 0xFF, 0xCC), bold=True,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, cx + Inches(0.1), cy + Inches(0.33), cw - Inches(0.15),
                    Inches(0.38), title, font_size=13, font_color=C_WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

        item_y = cy + Inches(0.9)
        for item in items:
            tb = slide.shapes.add_textbox(cx + Inches(0.15), item_y,
                                          cw - Inches(0.25), Inches(0.9))
            tb.word_wrap = True
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = f"• {item}"; r.font.size = Pt(11)
            r.font.color.rgb = C_BLACK; r.font.name = "Apple SD Gothic Neo"
            item_y += Inches(0.95)


# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Building slides...")

    slide_01_cover(prs)
    print("  [1/14] Cover")

    slide_02_current_state(prs)
    print("  [2/14] 현황 & 과제")

    slide_03_before_after(prs)
    print("  [3/14] Before vs After")

    slide_04_writing_process(prs)
    print("  [4/14] TC 작성 프로세스")

    slide_05_management_structure(prs)
    print("  [5/14] TC 관리 구조")

    slide_06_automation_process(prs)
    print("  [6/14] 자동화 개발 프로세스")

    slide_07_execution(prs)
    print("  [7/14] 실행 & 결과 관리")

    slide_08_dashboard(prs)
    print("  [8/14] 대시보드")

    slide_09_requirements(prs)
    print("  [9/14] 시스템 요구사항")

    slide_10_tech_stack(prs)
    print("  [10/14] 기술 스택")

    slide_11_data_structure(prs)
    print("  [11/14] 데이터 구조")

    slide_12_change_management(prs)
    print("  [12/14] 변경 관리")

    slide_13_conclusion(prs)
    print("  [13/14] 결론")

    slide_14_next_steps(prs)
    print("  [14/14] Next Steps")

    today = datetime.date.today().strftime("%Y%m%d")
    out_dir = "/Users/myungjookim/_claude26/pptx"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"SPCY_TC_Proposal_Presentation_{today}.pptx")

    prs.save(out_path)
    print(f"\n✅  저장 완료: {out_path}")
    return out_path


if __name__ == "__main__":
    build()
