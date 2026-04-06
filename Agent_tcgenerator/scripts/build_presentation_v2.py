#!/usr/bin/env python3
"""
TC Management System Presentation Builder v2
- Slide 6 (blank) removed
- All diagram slides enhanced with real connectors, callouts, badges
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os
import datetime

# ── Color Palette ──────────────────────────────────────
C_NAVY   = RGBColor(0x1B, 0x2E, 0x4B)
C_BLUE   = RGBColor(0x23, 0x6F, 0xAB)
C_TEAL   = RGBColor(0x00, 0xA8, 0xA8)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY   = RGBColor(0x74, 0x7D, 0x8C)
C_LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)
C_YELLOW = RGBColor(0xF3, 0x9C, 0x12)
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
C_DIVIDER= RGBColor(0xD5, 0xDB, 0xE5)

W = Inches(10)
H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.5)
FONT = "Apple SD Gothic Neo"


# ── Primitives ──────────────────────────────────────────

def hex6(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def add_rect(slide, x, y, w, h, fill=None, line=None, lw_pt=None, radius=None):
    """Add rectangle (optionally rounded)."""
    if radius:
        shape = slide.shapes.add_shape(5, x, y, w, h)  # rounded rect
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if lw_pt:
            shape.line.width = Pt(lw_pt)
    else:
        shape.line.fill.background()
    return shape


def add_tb(slide, x, y, w, h, text, size=16, color=C_BLACK,
           bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    return tb


def add_multiline(slide, x, y, w, h, lines, default_size=14, wrap=True):
    """lines = list of (text, size, color, bold, align) tuples or strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, color, bold, align = item, default_size, C_BLACK, False, PP_ALIGN.LEFT
        elif len(item) == 2:
            text, size = item; color, bold, align = C_BLACK, False, PP_ALIGN.LEFT
        elif len(item) == 3:
            text, size, color = item; bold, align = False, PP_ALIGN.LEFT
        elif len(item) == 4:
            text, size, color, bold = item; align = PP_ALIGN.LEFT
        else:
            text, size, color, bold, align = item

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = FONT
    return tb


def add_connector(slide, x1, y1, x2, y2, color=C_GRAY, width_pt=1.5, arrow=True):
    """Straight connector with optional arrowhead at end."""
    conn = slide.shapes.add_connector(1, int(x1), int(y1), int(x2), int(y2))
    cxnSp = conn._element
    spPr = cxnSp.find(qn('p:spPr'))
    for old in spPr.findall(qn('a:ln')):
        spPr.remove(old)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(int(Pt(width_pt))))
    sf = etree.SubElement(ln, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr'))
    sc.set('val', hex6(color))
    if arrow:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'arrow')
        tail.set('w', 'med')
        tail.set('len', 'med')
    return conn


def add_badge(slide, x, y, text, bg=C_BLUE, fg=C_WHITE, size=11):
    """Small colored pill badge."""
    w = Inches(1.1)
    h = Inches(0.3)
    add_rect(slide, x, y, w, h, fill=bg, radius=0.12)
    add_tb(slide, x, y + Inches(0.01), w, h - Inches(0.02),
           text, size=size, color=fg, bold=True, align=PP_ALIGN.CENTER)


def add_callout(slide, x, y, w, h, text, bg=C_LGRAY, fg=C_NAVY,
                size=12, border=C_BLUE, lw=1.0):
    """Rounded callout box."""
    add_rect(slide, x, y, w, h, fill=bg, line=border, lw_pt=lw, radius=0.08)
    add_tb(slide, x + Inches(0.1), y + Inches(0.05),
           w - Inches(0.2), h - Inches(0.1),
           text, size=size, color=fg, bold=True, wrap=True)


def add_title_bar(slide, title, sub=None, bg=C_NAVY):
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, W, bar_h, fill=bg)
    add_rect(slide, 0, bar_h - Inches(0.04), W, Inches(0.04), fill=C_TEAL)
    add_tb(slide, MARGIN, Inches(0.2), W - MARGIN * 2, Inches(0.7),
           title, size=28, color=C_WHITE, bold=True)
    if sub:
        add_tb(slide, MARGIN, Inches(0.78), W - MARGIN * 2, Inches(0.35),
               sub, size=12, color=C_TEAL)


# ══════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════

def s01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_NAVY)
    add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), fill=C_TEAL)
    add_rect(slide, 0, 0, Inches(0.35), H, fill=C_BLUE)
    add_tb(slide, Inches(0.7), Inches(1.5), Inches(8), Inches(0.45),
           "QA TEAM  ·  TC MANAGEMENT SYSTEM PROPOSAL",
           size=12, color=C_TEAL, bold=True)
    add_tb(slide, Inches(0.7), Inches(2.1), Inches(8.5), Inches(1.6),
           "테스트 케이스 관리 시스템 제안",
           size=40, color=C_WHITE, bold=True)
    add_tb(slide, Inches(0.7), Inches(3.9), Inches(8), Inches(0.55),
           "TC 작성 · 관리 · 실행 · 결과까지 — 단일 통합 시스템 설계",
           size=18, color=RGBColor(0xA8, 0xD8, 0xEA))
    add_rect(slide, Inches(0.7), Inches(4.55), Inches(4), Inches(0.04), fill=C_TEAL)
    today = datetime.date.today().strftime("%Y.%m.%d")
    add_tb(slide, Inches(0.7), Inches(4.75), Inches(8), Inches(0.4),
           f"v0.2   |   {today}", size=13, color=C_GRAY)
    for i in range(5):
        for j in range(5):
            a = 0x20 + (i + j) * 0x08
            c = RGBColor(min(a, 0xFF), min(a + 0x30, 0xFF), min(a + 0x60, 0xFF))
            add_rect(slide, Inches(8.0 + j * 0.25), Inches(0.5 + i * 0.25),
                     Inches(0.18), Inches(0.18), fill=c)


def s02_current_state(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "TC 자동화 개발 — 현황 & 과제")

    for bx, bg, title_c, hdr_txt, items in [
        (MARGIN, C_BLUE, C_BLUE, "📋  현황",
         [("전체 TC: v2.0 기준 정리 완료", 14, C_BLACK, False),
          ("→ v3.0에 자동화 여부 기록 완료", 13, C_GRAY, False),
          ("", 6, C_BLACK, False),
          ("Phase 구분", 14, C_BLUE, True),
          ("  P1: Youthmeta (5개 거래소)", 13, C_BLACK, False),
          ("  P2: Mobile (iOS / Android)", 13, C_BLACK, False),
          ("", 6, C_BLACK, False),
          ("TC 형식: BDD 기반", 14, C_BLACK, False),
          ("Excel + Markdown 병행 관리", 13, C_GRAY, False)]),
        (Inches(5.1), C_ORANGE, C_ORANGE, "⚠️  필요한 작업",
         [("① TC ID 체계 정비", 14, C_NAVY, True),
          ("   SPCY-{Phase}-{번호} 전체 소급", 13, C_GRAY, False),
          ("", 6, C_BLACK, False),
          ("② TC 추가 작성", 14, C_NAVY, True),
          ("   Youthmeta / Mobile 신규 작성", 13, C_GRAY, False),
          ("", 6, C_BLACK, False),
          ("③ Phase 1 TC 리뷰", 14, C_NAVY, True),
          ("   거래소별 커버리지 검증", 13, C_GRAY, False),
          ("", 6, C_BLACK, False),
          ("④ 자동화 연동", 14, C_NAVY, True),
          ("   is_automated / auto_tc_id 추가", 13, C_GRAY, False)])
    ]:
        bw, bh = Inches(4.2), Inches(5.1)
        add_rect(slide, bx, CONTENT_TOP, bw, bh, fill=C_WHITE)
        add_rect(slide, bx, CONTENT_TOP, bw, Inches(0.45), fill=bg)
        add_tb(slide, bx + Inches(0.2), CONTENT_TOP + Inches(0.05),
               bw - Inches(0.3), Inches(0.38),
               hdr_txt, size=15, color=C_WHITE, bold=True)
        add_multiline(slide, bx + Inches(0.2), CONTENT_TOP + Inches(0.55),
                      bw - Inches(0.35), bh - Inches(0.65), items)


def s03_before_after(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "Before vs After — 시스템 도입 효과")

    col_w = Inches(4.35)
    gap = Inches(0.2)
    top_y = CONTENT_TOP
    col_h = Inches(5.15)

    rows = [
        ("TC 작성",      "수동 (QA 개인 역량 의존)",              "Claude + Harness 자동 초안"),
        ("TC 형식",      "작성자마다 다름",                       "BDD / 규칙 파일 기준 통일"),
        ("TC 저장",      "개별 Excel, 네이밍 비일관",             "SPCY_TC_{Phase}_{날짜}_v{N}"),
        ("TC ID",       "없거나 파일마다 달라 추적 불가",          "SPCY-{Phase}-{번호} 고유 식별"),
        ("자동화 연동",   "TC와 자동화 코드 분리",                  "is_automated + auto_tc_id"),
        ("결과 기록",    "실행 후 별도 업데이트",                  "수동 + 자동 결과 자동 수집"),
        ("현황 파악",    "파일 열어봐야 알 수 있음",               "대시보드 실시간 진행률 확인"),
        ("리포트",       "별도 수작업 작성",                      "Test Result Agent 자동 생성"),
    ]

    for col, (clr, hdr) in enumerate([(RGBColor(0xB7,0x3A,0x3A), "✗  Before"),
                                       (C_GREEN, "✓  After")]):
        x = MARGIN + col * (col_w + gap)
        add_rect(slide, x, top_y, col_w, col_h, fill=C_WHITE)
        add_rect(slide, x, top_y, col_w, Inches(0.45), fill=clr)
        add_tb(slide, x + Inches(0.2), top_y + Inches(0.05),
               col_w - Inches(0.3), Inches(0.38),
               hdr, size=16, color=C_WHITE, bold=True)
        cy = top_y + Inches(0.52)
        for label, before, after in rows:
            text = before if col == 0 else after
            fg = RGBColor(0x55,0x55,0x55) if col == 0 else RGBColor(0x22,0x66,0x33)
            add_rect(slide, x + Inches(0.18), cy + Inches(0.05),
                     Inches(0.06), Inches(0.38), fill=clr)
            add_tb(slide, x + Inches(0.3), cy, col_w - Inches(0.42),
                   Inches(0.18), label, size=10, color=C_GRAY, bold=True)
            add_tb(slide, x + Inches(0.3), cy + Inches(0.17), col_w - Inches(0.42),
                   Inches(0.35), text, size=12, color=fg)
            cy += Inches(0.57)

    # VS badge center
    vx = MARGIN + col_w + gap/2 - Inches(0.2)
    vy = top_y + col_h/2 - Inches(0.32)
    add_rect(slide, vx, vy, Inches(0.4), Inches(0.64), fill=C_NAVY, radius=0.05)
    add_tb(slide, vx, vy + Inches(0.1), Inches(0.4), Inches(0.44),
           "VS", size=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


def s04_writing_process(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "TC 작성 프로세스")

    steps = [
        (C_BLUE,   "①", "초기 작성",   "TC Generator\n(Claude + Harness)"),
        (C_TEAL,   "②", "QA 리뷰",    "도메인별 커버리지\n검증 및 보강"),
        (C_GREEN,  "③", "수정 반영",   "누락 기능\n자동 재작성"),
        (C_ORANGE, "④", "산출물 배포", "Excel + Markdown\n버전 관리 출력"),
        (C_NAVY,   "⑤", "버전 관리",  "Phase / Sprint 단위\nMinor 버전 증가"),
    ]

    sw = Inches(1.6)
    sh = Inches(3.2)
    gap = Inches(0.15)
    sy = Inches(1.9)
    total_w = len(steps) * sw + (len(steps)-1) * gap
    sx = (W - total_w) / 2

    for i, (color, num, title, desc) in enumerate(steps):
        x = sx + i * (sw + gap)
        add_rect(slide, x, sy, sw, sh, fill=color)

        # number circle
        add_rect(slide, x + sw/2 - Inches(0.28), sy + Inches(0.18),
                 Inches(0.56), Inches(0.56), fill=C_WHITE, radius=0.5)
        add_tb(slide, x + sw/2 - Inches(0.28), sy + Inches(0.15),
               Inches(0.56), Inches(0.56),
               num, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)

        add_tb(slide, x + Inches(0.1), sy + Inches(0.9), sw - Inches(0.15),
               Inches(0.45), title, size=14, color=C_WHITE, bold=True,
               align=PP_ALIGN.CENTER)

        tb = slide.shapes.add_textbox(x + Inches(0.1), sy + Inches(1.45),
                                      sw - Inches(0.15), Inches(1.5))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = desc
        r.font.size = Pt(12); r.font.name = FONT
        r.font.color.rgb = RGBColor(0xD8, 0xEC, 0xFF)

        # real connector arrow to next step
        if i < len(steps) - 1:
            ax = x + sw + Inches(0.01)
            ay = sy + sh / 2
            bx2 = x + sw + gap - Inches(0.01)
            add_connector(slide, ax, ay, bx2, ay, color=C_GRAY, width_pt=1.5)

    # Output label
    add_rect(slide, sx, sy + sh + Inches(0.18), total_w, Inches(0.55),
             fill=C_WHITE, line=C_DIVIDER if hasattr(RGBColor, '__iter__') else None,
             lw_pt=0.5)
    add_tb(slide, sx + Inches(0.3), sy + sh + Inches(0.22), total_w - Inches(0.5),
           Inches(0.4),
           "결과물: SPCY_TC_{Phase}_{YYYYMMDD}_v{N}.xlsx  +  .md",
           size=12, color=C_GRAY, align=PP_ALIGN.CENTER)


def s05_management_structure(prs):
    """TC 관리 구조 — with side callouts, arrow connectors, TC ID badges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "TC 관리 구조")

    # ── Master Repository (top center) ──
    mw, mh = Inches(4.2), Inches(0.95)
    mx = (W - mw) / 2
    my = CONTENT_TOP + Inches(0.05)
    add_rect(slide, mx, my, mw, mh, fill=C_NAVY)
    add_rect(slide, mx, my, mw, Inches(0.04), fill=C_TEAL)
    add_tb(slide, mx + Inches(0.25), my + Inches(0.06), mw - Inches(0.4),
           Inches(0.38), "📦  Master TC Repository",
           size=16, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, mx + Inches(0.1), my + Inches(0.52), mw - Inches(0.2),
           Inches(0.35),
           "전체 TC 통합 관리  |  TC ID: UNIQUE  |  버전 관리",
           size=11, color=C_TEAL, align=PP_ALIGN.CENTER)

    # ── Left callout ──
    add_callout(slide, Inches(0.3), my + Inches(0.15), Inches(1.8), Inches(0.65),
                "품  전체 TC\n통합 관리", bg=C_WHITE, fg=C_NAVY, border=C_BLUE, size=11)
    add_connector(slide, Inches(2.1), my + Inches(0.47),
                  mx, my + Inches(0.47), color=C_BLUE, width_pt=1.2, arrow=False)

    # ── Right callout ──
    add_callout(slide, Inches(7.9), my + Inches(0.15), Inches(1.8), Inches(0.65),
                "⏱  버전\n히스토리 보존", bg=C_WHITE, fg=C_NAVY, border=C_TEAL, size=11)
    add_connector(slide, mx + mw, my + Inches(0.47),
                  Inches(7.9), my + Inches(0.47), color=C_TEAL, width_pt=1.2, arrow=False)

    # ── Center connector label ──
    cl_w, cl_h = Inches(3.4), Inches(0.65)
    cl_x = (W - cl_w) / 2
    cl_y = my + mh + Inches(0.12)
    add_rect(slide, cl_x, cl_y, cl_w, cl_h, fill=C_WHITE,
             line=C_GRAY, lw_pt=0.8, radius=0.06)
    add_tb(slide, cl_x + Inches(0.15), cl_y + Inches(0.03),
           cl_w - Inches(0.3), Inches(0.3),
           "고유 ID 발급 및 체계 정비",
           size=12, color=C_NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, cl_x + Inches(0.15), cl_y + Inches(0.33),
           cl_w - Inches(0.3), Inches(0.28),
           "SPCY-{Phase}-{번호}",
           size=11, color=C_BLUE, align=PP_ALIGN.CENTER)

    # vertical connector: master → label
    add_connector(slide, W/2, my + mh, W/2, cl_y,
                  color=C_GRAY, width_pt=1.5, arrow=False)
    # vertical connector: label → phase row
    ph_y = cl_y + cl_h + Inches(0.12)
    add_connector(slide, W/2, cl_y + cl_h, W/2, ph_y,
                  color=C_GRAY, width_pt=1.5, arrow=False)

    # ── Phase boxes ──
    phases = [
        (C_BLUE,  "Phase 1\nYouthmeta",  "Gate · OKX · BYBIT\nBitget · Hyperliquid", "P1-010"),
        (C_TEAL,  "Phase 2\nMobile",     "iOS · Android",                            "P2-105"),
        (C_GREEN, "Phase N\n신규 추가",   "확장 가능 구조\nPN-{번호}",               "PN-XXX"),
    ]
    pw, ph_h = Inches(2.6), Inches(2.0)
    gap = Inches(0.35)
    total_pw = len(phases) * pw + (len(phases)-1) * gap
    px_start = (W - total_pw) / 2

    for i, (color, title, detail, badge_id) in enumerate(phases):
        px = px_start + i * (pw + gap)

        add_rect(slide, px, ph_y, pw, ph_h, fill=color)
        # Gradient-like header bar
        add_rect(slide, px, ph_y, pw, Inches(0.08), fill=C_WHITE)

        # TC ID badge (top-right corner)
        add_rect(slide, px + pw - Inches(1.05), ph_y + Inches(0.1),
                 Inches(0.95), Inches(0.28), fill=C_WHITE, radius=0.1)
        add_tb(slide, px + pw - Inches(1.05), ph_y + Inches(0.1),
               Inches(0.95), Inches(0.28),
               badge_id, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)

        add_tb(slide, px + Inches(0.12), ph_y + Inches(0.15), pw - Inches(0.15),
               Inches(0.6), title, size=14, color=C_WHITE, bold=True,
               align=PP_ALIGN.CENTER)

        add_tb(slide, px + Inches(0.12), ph_y + Inches(0.82), pw - Inches(0.15),
               Inches(0.9), detail, size=12,
               color=RGBColor(0xD0, 0xF0, 0xFF), align=PP_ALIGN.CENTER)

        # connector from label/center to this phase box
        phase_cx = px + pw / 2
        add_connector(slide, W/2, ph_y, phase_cx, ph_y,
                      color=color, width_pt=1.5, arrow=False)
        add_connector(slide, phase_cx, ph_y - Inches(0.01), phase_cx, ph_y + Inches(0.02),
                      color=color, width_pt=1.5, arrow=True)

    # ── Bottom principles bar ──
    bar_y = ph_y + ph_h + Inches(0.2)
    bar_h = Inches(0.6)
    add_rect(slide, MARGIN, bar_y, W - MARGIN*2, bar_h, fill=C_NAVY)
    add_tb(slide, MARGIN + Inches(0.2), bar_y + Inches(0.06),
           Inches(2.0), Inches(0.35),
           "TC 관리 핵심 원칙", size=12, color=C_TEAL, bold=True)
    items = ["TC ID는 고유 식별자", "Phase별 독립 관리", "자동화 여부 포함", "버전 히스토리 보존"]
    colors_b = [C_BLUE, C_TEAL, C_GREEN, C_ORANGE]
    for j, (item, c) in enumerate(zip(items, colors_b)):
        x = MARGIN + Inches(2.3) + j * Inches(1.85)
        add_rect(slide, x - Inches(0.1), bar_y + Inches(0.14),
                 Inches(0.06), Inches(0.3), fill=c)
        add_tb(slide, x, bar_y + Inches(0.14), Inches(1.7), Inches(0.3),
               item, size=11, color=C_WHITE)


def s06_automation_process(prs):
    """TC 자동화 개발 프로세스 — enhanced step flow with role transitions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "TC 자동화 개발 프로세스")

    steps = [
        (C_BLUE,   "QA",        "① Phase별 TC 초안 작성",
         "Claude + Harness 자동 생성"),
        (C_TEAL,   "QA → 개발", "② TC 전달 + 자동화 여부 공유",
         "TC ID 기반으로 자동화 대상 명시"),
        (C_GREEN,  "개발",      "③ 자동화 케이스 확정",
         "구현 가능 여부 판단"),
        (C_ORANGE, "QA",        "④ 자동 / 수동 분류 및 스케줄 확정",
         "실행 계획 수립"),
        (C_NAVY,   "QA + 개발", "⑤ 테스트 실행 ↔ 자동화 구현",
         "병렬 진행, 결과 통합"),
    ]

    step_h = Inches(0.82)
    sy = CONTENT_TOP
    left_col_w = Inches(1.2)
    step_x = MARGIN + left_col_w + Inches(0.15)
    step_w = W - step_x - MARGIN

    for i, (color, role, title, desc) in enumerate(steps):
        y = sy + i * (step_h + Inches(0.1))

        # role badge (left column)
        add_rect(slide, MARGIN, y, left_col_w, step_h, fill=color)
        add_tb(slide, MARGIN, y + Inches(0.22), left_col_w, Inches(0.38),
               role, size=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

        # main step box
        add_rect(slide, step_x, y, step_w, step_h, fill=C_WHITE)
        add_rect(slide, step_x, y, Inches(0.06), step_h, fill=color)

        add_tb(slide, step_x + Inches(0.2), y + Inches(0.08), Inches(5.0), Inches(0.38),
               title, size=15, color=color, bold=True)
        add_tb(slide, step_x + Inches(0.2), y + Inches(0.46), step_w - Inches(0.3),
               Inches(0.3), desc, size=12, color=C_GRAY)

        # connector arrow between steps
        if i < len(steps) - 1:
            ay = y + step_h + Inches(0.02)
            by2 = ay + Inches(0.06)
            # Left column connector
            add_connector(slide, MARGIN + left_col_w/2, ay,
                          MARGIN + left_col_w/2, by2,
                          color=C_GRAY, width_pt=1.2, arrow=True)

    # right side note
    note_x = W - MARGIN - Inches(2.5)
    note_y = sy + Inches(0.2)
    add_rect(slide, note_x, note_y, Inches(2.5), Inches(4.5),
             fill=RGBColor(0xE8, 0xF4, 0xFD), line=C_BLUE, lw_pt=0.5)
    add_tb(slide, note_x + Inches(0.15), note_y + Inches(0.1),
           Inches(2.2), Inches(0.35),
           "핵심 원칙", size=12, color=C_BLUE, bold=True)
    notes_text = [
        "• TC ID 기반 협업",
        "• 자동화 여부 사전 확정",
        "• QA / 개발 병렬 진행",
        "• 결과 통합 관리",
        "• 커버리지 ≥ 95%",
    ]
    cy = note_y + Inches(0.5)
    for note in notes_text:
        add_tb(slide, note_x + Inches(0.15), cy, Inches(2.2), Inches(0.4),
               note, size=11, color=C_NAVY)
        cy += Inches(0.42)


def s07_execution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "TC 실행 & 결과 관리")

    col_w = Inches(4.3)
    gap = Inches(0.3)

    for col, (bg, hdr, items) in enumerate([
        (C_BLUE, "▶  TC 실행 프로세스", [
            (C_BLUE,   "수동 TC",   "QA에서 버전별 Excel에 결과 직접 기록"),
            (C_TEAL,   "자동 TC",   "자동화 시스템이 결과 자동 기록\nis_automated = true"),
            (C_ORANGE, "실행 중 변경", "TC 추가 / 수정 / 삭제 가능\nMinor 버전 올려 관리"),
        ]),
        (C_TEAL, "📊  결과 관리", [
            (C_NAVY,  "목표 일정 관리",  "시험 기간 설정 및 진행률 추적"),
            (C_BLUE,  "실시간 모니터링", "Pass / Fail / N/T 현황 대시보드"),
            (C_TEAL,  "자동화팀 역할",  "이슈 디버깅 & 케이스 보고"),
            (C_GREEN, "결과서 자동 생성", "Test Result Agent 기반\nSlack / 메일 자동 발송"),
        ]),
    ]):
        cx = MARGIN + col * (col_w + gap)
        ch = Inches(5.2)
        add_rect(slide, cx, CONTENT_TOP, col_w, ch, fill=C_WHITE)
        add_rect(slide, cx, CONTENT_TOP, col_w, Inches(0.45), fill=bg)
        add_tb(slide, cx + Inches(0.18), CONTENT_TOP + Inches(0.05),
               col_w - Inches(0.3), Inches(0.38),
               hdr, size=14, color=C_WHITE, bold=True)

        iy = CONTENT_TOP + Inches(0.55)
        for color, label, detail in items:
            h_item = Inches(1.05)
            add_rect(slide, cx + Inches(0.18), iy, Inches(0.06), h_item - Inches(0.1),
                     fill=color)
            add_tb(slide, cx + Inches(0.35), iy + Inches(0.02),
                   col_w - Inches(0.5), Inches(0.32),
                   label, size=13, color=color, bold=True)
            tb = slide.shapes.add_textbox(cx + Inches(0.35), iy + Inches(0.32),
                                          col_w - Inches(0.5), Inches(0.62))
            tb.word_wrap = True
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = detail; r.font.size = Pt(11); r.font.color.rgb = C_GRAY
            r.font.name = FONT
            iy += h_item

    # status bar at bottom
    statuses = [("Pass", C_GREEN), ("Fail", C_RED), ("N/T", C_GRAY), ("N/A", RGBColor(0xBD,0xBD,0xBD))]
    bar_y = CONTENT_TOP + Inches(5.3)
    bw = (W - MARGIN*2 - Inches(0.3)*3) / 4
    for j, (label, c) in enumerate(statuses):
        bx = MARGIN + j * (bw + Inches(0.3))
        add_rect(slide, bx, bar_y, bw, Inches(0.38), fill=c, radius=0.05)
        add_tb(slide, bx, bar_y + Inches(0.04), bw, Inches(0.3),
               label, size=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


def s08_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "대시보드 기능")

    cards = [
        (C_RED,    "🔴  Fail 리스트",
         "중요도별 목록 (High → Low)\n자동 리포트 → 메일 / Slack\n\n예: Critical 3건 / Major 12건"),
        (C_BLUE,   "🔍  검색 & 필터",
         "거래소별 / Phase별 결과 필터\n키워드 전체 검색\n\n예: \"로그인\" → 해당 TC 목록"),
        (C_GREEN,  "📈  통계 & 진행률",
         "Pass / Fail / N/T 실시간 현황\n버전별 비교 차트\n\n예: v3.0 Pass율 84% / 전버전比 +6%"),
        (C_ORANGE, "🤖  자동화 비율",
         "전체 대비 자동화 TC 비율\nPhase별 / 도메인별 자동화 현황\n\n예: P2 Mobile 자동화 37%"),
    ]

    cw = Inches(4.3)
    ch = Inches(2.35)
    gap = Inches(0.3)
    for i, (color, title, detail) in enumerate(cards):
        row, col = divmod(i, 2)
        cx = MARGIN + col * (cw + gap)
        cy = CONTENT_TOP + row * (ch + gap)
        add_rect(slide, cx, cy, cw, ch, fill=C_WHITE)
        add_rect(slide, cx, cy, Inches(0.1), ch, fill=color)
        add_tb(slide, cx + Inches(0.2), cy + Inches(0.12), cw - Inches(0.3),
               Inches(0.4), title, size=14, color=color, bold=True)
        add_rect(slide, cx + Inches(0.2), cy + Inches(0.54), cw - Inches(0.35),
                 Inches(0.02), fill=RGBColor(0xE0,0xE0,0xE0))
        tb = slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.62),
                                      cw - Inches(0.3), ch - Inches(0.75))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = detail; r.font.size = Pt(12); r.font.color.rgb = C_GRAY
        r.font.name = FONT


def s09_requirements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "시스템 요구사항")

    reqs = [
        (C_BLUE,   "👁  Viewer 기능",
         "컬럼별 선택 옵션으로 TC 보기\nExcel / CSV export 지원\n거래소 / Phase / 도메인 필터\n\n→ 팀원 누구나 TC 현황 파악 가능"),
        (C_TEAL,   "✏️  Edit 기능",
         "Edit → 임시저장 → 승인 → TC 업데이트\n승인 단계 Skip 옵션 지원\n변경 히스토리 자동 기록\n\n→ 실수 없는 TC 수정 플로우"),
        (C_GREEN,  "👥  협업 환경",
         "1인 또는 팀 작업 모드\n동시 편집 / 권한 관리 지원\n테스터 배정 (Tester 1 / 2)\n\n→ 병렬 테스트 실행 지원"),
        (C_ORANGE, "🔗  자동화 연동",
         "자동화 시스템과 TC ID 연동\nis_automated / auto_tc_id 필드\n결과 자동 수집 API\n\n→ 수동/자동 결과 통합 관리"),
    ]

    cw = Inches(4.3)
    ch = Inches(2.35)
    gap = Inches(0.3)
    for i, (color, title, detail) in enumerate(reqs):
        row, col = divmod(i, 2)
        cx = MARGIN + col * (cw + gap)
        cy = CONTENT_TOP + row * (ch + gap)
        add_rect(slide, cx, cy, cw, ch, fill=C_WHITE)
        add_rect(slide, cx, cy, Inches(0.1), ch, fill=color)
        add_tb(slide, cx + Inches(0.2), cy + Inches(0.12), cw - Inches(0.3),
               Inches(0.4), title, size=14, color=color, bold=True)
        add_rect(slide, cx + Inches(0.2), cy + Inches(0.54), cw - Inches(0.35),
                 Inches(0.02), fill=RGBColor(0xE0,0xE0,0xE0))
        tb = slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.62),
                                      cw - Inches(0.3), ch - Inches(0.75))
        tb.word_wrap = True
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = detail; r.font.size = Pt(11); r.font.color.rgb = C_GRAY
        r.font.name = FONT


def s10_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "기술 스택")

    # Flow arrow between columns
    stacks = [
        (C_NAVY,   "TC 생성",
         ["Claude Sonnet (AI 모델)",
          "Harness 오케스트레이터",
          "파이프라인:",
          "  doc-parser → tc-writer",
          "  → tc-reviewer",
          "규칙 파일 기반 통제"]),
        (C_BLUE,   "TC 저장 / 배포",
         ["Excel (.xlsx) + MD (.md)",
          "파일명 규칙:",
          "  SPCY_TC_{Phase}",
          "  _{날짜}_v{N}",
          "build_excel.py (openpyxl)",
          "버전 자동 증가"]),
        (C_TEAL,   "TC 조회 / 실행",
         ["Viewer: 웹 기반 (예정)",
          "컬럼 선택 / 필터 / Export",
          "수동: Excel 직접 기록",
          "자동: is_automated 연동",
          "auto_tc_id 매핑",
          "결과 자동 수집"]),
        (C_GREEN,  "결과 리포팅",
         ["Test Result Agent",
          "결과서 자동 생성",
          "Slack / 메일 리포트",
          "Fail 중요도별 분류",
          "진행률 대시보드",
          "버전별 비교"]),
    ]

    sw = Inches(1.95)
    sh = Inches(5.0)
    gap = Inches(0.28)
    sx = MARGIN + Inches(0.15)

    for i, (color, title, items) in enumerate(stacks):
        x = sx + i * (sw + gap)
        add_rect(slide, x, CONTENT_TOP, sw, sh, fill=C_WHITE)
        add_rect(slide, x, CONTENT_TOP, sw, Inches(0.48), fill=color)
        add_tb(slide, x + Inches(0.1), CONTENT_TOP + Inches(0.07),
               sw - Inches(0.15), Inches(0.38),
               title, size=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

        cy = CONTENT_TOP + Inches(0.6)
        for item in items:
            is_sub = item.startswith("  ")
            c = C_GRAY if is_sub else C_BLACK
            tb = slide.shapes.add_textbox(x + Inches(0.12), cy, sw - Inches(0.2), Inches(0.42))
            tb.word_wrap = True
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = ("  " + item.strip() if is_sub else f"• {item}")
            r.font.size = Pt(10 if is_sub else 11)
            r.font.color.rgb = c; r.font.name = FONT
            cy += Inches(0.44 if is_sub else 0.5)

        # arrow to next
        if i < len(stacks) - 1:
            ax = x + sw + Inches(0.02)
            ay = CONTENT_TOP + sh/2
            add_connector(slide, ax, ay, ax + gap - Inches(0.04), ay,
                          color=color, width_pt=1.5)


def s11_data_structure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "TC 데이터 구조 설계")

    groups = [
        (C_NAVY,   "기본 식별",
         "tc_id  |  phase  |  domain  |  version",
         "예: SPCY-P2-0042  |  P2_Mobile  |  AUTH  |  v3.1"),
        (C_BLUE,   "테스트 내용",
         "title · priority(H/M/L) · type(Pos/Neg/Edge)",
         "precondition / given · when · then · expected_result / platform"),
        (C_TEAL,   "거래소 필드  (Phase 1 전용)",
         "Gate · OKX · Bybit · Bitget · Hyperliquid",
         "→ Pass / Fail / N/T / N/A   |   tester: Tester1 / Tester2"),
        (C_ORANGE, "자동화 연동",
         "is_automated: true/false   |   auto_tc_id: 자동화 시스템 ID",
         "automation_status: 완료 / 진행중 / 미정"),
        (C_GREEN,  "결과 & 이력",
         "status: Pass/Fail/N/T/N/A   |   actual_result   |   tested_at   |   tester",
         "change_log: [ { version, who, what, when } ]"),
    ]

    gh = Inches(0.8)
    gap = Inches(0.09)
    gy = CONTENT_TOP

    for color, title, fields, example in groups:
        # Left accent bar
        add_rect(slide, MARGIN, gy, Inches(0.1), gh, fill=color)
        # White content area
        add_rect(slide, MARGIN + Inches(0.14), gy, W - MARGIN*2 - Inches(0.14), gh,
                 fill=C_WHITE)
        # Title
        add_tb(slide, MARGIN + Inches(0.28), gy + Inches(0.06),
               Inches(1.85), Inches(0.3),
               title, size=12, color=color, bold=True)
        # Fields (code style)
        add_tb(slide, MARGIN + Inches(2.2), gy + Inches(0.05),
               Inches(7.0), Inches(0.3),
               fields, size=11, color=C_BLACK)
        # Example (lighter)
        add_tb(slide, MARGIN + Inches(2.2), gy + Inches(0.42),
               Inches(7.0), Inches(0.28),
               example, size=10, color=C_GRAY, italic=True)
        gy += gh + gap


def s12_change_management(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "변경 관리")

    phases_data = [
        (C_BLUE,   "테스트 시작 전",
         [("TC 수정 시 자동화팀 공유", C_NAVY, True),
          ("→ 자동화 여부 컨펌 프로세스 적용", C_BLACK, False),
          ("→ 전체 TC 반영 (수시 또는 일괄)", C_BLACK, False),
          ("", C_BLACK, False),
          ("타이밍: 수시 또는 sprint 주기", C_GRAY, False)]),
        (C_ORANGE, "테스트 진행 중",
         [("기본 원칙:", C_NAVY, True),
          ("Minor 버전 올린 후 tc id 재관리", C_BLACK, False),
          ("스텝 / 기대결과 변경 내역 기록", C_BLACK, False),
          ("", C_BLACK, False),
          ("긴급 수정:", C_RED, True),
          ("즉시 공유 → 반영 → 재진행", C_BLACK, False)]),
        (C_GREEN,  "변경 히스토리",
         [("tc id + 수정 내용 기록 (Optional)", C_NAVY, True),
          ("변경 날짜 & 작성자 보존", C_BLACK, False),
          ("change_log 필드로 관리", C_BLACK, False),
          ("", C_BLACK, False),
          ("형식: {version, who, what, when}", C_GRAY, False)]),
    ]

    # horizontal timeline connector
    cw = Inches(2.75)
    ch = Inches(4.6)
    gap = Inches(0.35)
    total_w = len(phases_data) * cw + (len(phases_data)-1) * gap
    cx_start = (W - total_w) / 2
    cy = CONTENT_TOP

    # Draw timeline line
    line_y = cy + Inches(0.24)
    add_rect(slide, cx_start + cw/2, line_y + Inches(0.02),
             total_w - cw, Inches(0.04), fill=C_GRAY)

    for i, (color, title, items) in enumerate(phases_data):
        cx = cx_start + i * (cw + gap)

        # timeline dot
        add_rect(slide, cx + cw/2 - Inches(0.14), line_y - Inches(0.1),
                 Inches(0.28), Inches(0.28), fill=color, radius=0.5)

        add_rect(slide, cx, cy + Inches(0.3), cw, ch, fill=C_WHITE)
        add_rect(slide, cx, cy + Inches(0.3), cw, Inches(0.48), fill=color)
        add_tb(slide, cx + Inches(0.15), cy + Inches(0.35),
               cw - Inches(0.25), Inches(0.38),
               title, size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

        iy = cy + Inches(0.88)
        for text, fg, bold in items:
            if text == "":
                iy += Inches(0.15)
                continue
            tb = slide.shapes.add_textbox(cx + Inches(0.2), iy, cw - Inches(0.35),
                                          Inches(0.42))
            tb.word_wrap = True
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = text; r.font.size = Pt(12); r.font.color.rgb = fg
            r.font.bold = bold; r.font.name = FONT
            iy += Inches(0.42)


def s13_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_NAVY)
    add_rect(slide, 0, H - Inches(0.5), W, Inches(0.5), fill=C_TEAL)
    add_rect(slide, 0, 0, Inches(0.35), H, fill=C_BLUE)

    add_tb(slide, Inches(0.7), Inches(0.9), Inches(8), Inches(0.45),
           "CONCLUSION", size=12, color=C_TEAL, bold=True)
    add_tb(slide, Inches(0.7), Inches(1.3), Inches(8), Inches(0.65),
           "결론", size=36, color=C_WHITE, bold=True)
    add_rect(slide, Inches(0.7), Inches(2.0), Inches(3.5), Inches(0.04), fill=C_TEAL)

    points = [
        ("AI 기반 TC 자동 생성",   "Claude + Harness로 TC 초안 자동 생성 — QA 리소스 핵심 업무 집중",   C_BLUE),
        ("통합 관리 & 추적",       "Master TC + Phase별 관리로 누락·중복 없이 체계적 추적",              C_TEAL),
        ("자동화 팀 협업",         "개발팀과 TC ID 기반 협업 — 자동화 결과 자동 수집·통합",              C_GREEN),
        ("실시간 대시보드",         "진행률·Fail 현황 즉시 파악 — 리포트 자동 발송",                    C_ORANGE),
        ("QA 생산성 극대화",        "수동 작업 최소화, 단일 시스템으로 전 과정 통합",                   C_YELLOW),
    ]

    py = Inches(2.25)
    for label, desc, color in points:
        add_rect(slide, Inches(0.7), py + Inches(0.12), Inches(0.08), Inches(0.44), fill=color)
        add_tb(slide, Inches(0.92), py + Inches(0.08), Inches(2.5), Inches(0.32),
               label, size=14, color=color, bold=True)
        add_tb(slide, Inches(0.92), py + Inches(0.38), Inches(8.0), Inches(0.3),
               desc, size=12, color=RGBColor(0xCC, 0xD6, 0xE8))
        py += Inches(0.83)


def s14_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=C_LGRAY)
    add_title_bar(slide, "Next Steps  —  로드맵")

    phases = [
        (C_RED,    "즉시",   "TC 데이터 구조 확정",
         ["SPCY-{Phase}-{번호} 포맷 전체 소급",
          "Youthmeta / Mobile TC 마이그레이션",
          "자동화 필드 추가"]),
        (C_ORANGE, "단기",   "TC 관리 시스템 MVP",
         ["Viewer: 필터 + Export",
          "Edit: 임시저장 → 승인 플로우",
          "자동화 TC ID 연동"]),
        (C_BLUE,   "중기",   "실행 & 결과 연동",
         ["수동 / 자동 결과 기록 통합",
          "Fail 리스트 자동 리포트",
          "거래소별 현황 대시보드"]),
        (C_GREEN,  "장기",   "AI 결과 분석",
         ["Test Result Agent 자동 생성",
          "진행률 / 자동화 비율 통계",
          "회귀 테스트 자동 실행"]),
    ]

    cw = Inches(1.95)
    ch = Inches(4.8)
    gap = Inches(0.32)
    total_w = len(phases) * cw + (len(phases)-1) * gap
    cx_start = (W - total_w) / 2
    base_y = CONTENT_TOP + Inches(0.5)

    # Timeline bar
    tl_y = base_y - Inches(0.3)
    add_rect(slide, cx_start + cw/2, tl_y + Inches(0.12),
             total_w - cw, Inches(0.05), fill=RGBColor(0xBD,0xBD,0xBD))

    for i, (color, timing, title, items) in enumerate(phases):
        cx = cx_start + i * (cw + gap)

        # Timeline connector arrow between phases
        if i < len(phases) - 1:
            add_connector(slide, cx + cw, tl_y + Inches(0.14),
                          cx + cw + gap, tl_y + Inches(0.14),
                          color=color, width_pt=1.5)

        # Timeline dot
        add_rect(slide, cx + cw/2 - Inches(0.16), tl_y,
                 Inches(0.32), Inches(0.32), fill=color, radius=0.5)
        add_tb(slide, cx + cw/2 - Inches(0.6), tl_y - Inches(0.3), Inches(1.2),
               Inches(0.28), timing, size=12, color=color, bold=True,
               align=PP_ALIGN.CENTER)

        add_rect(slide, cx, base_y, cw, ch, fill=C_WHITE)
        add_rect(slide, cx, base_y, cw, Inches(0.6), fill=color)
        add_tb(slide, cx + Inches(0.1), base_y + Inches(0.06),
               cw - Inches(0.15), Inches(0.48),
               title, size=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

        iy = base_y + Inches(0.72)
        for item in items:
            tb = slide.shapes.add_textbox(cx + Inches(0.15), iy,
                                          cw - Inches(0.25), Inches(0.92))
            tb.word_wrap = True
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = f"• {item}"
            r.font.size = Pt(11); r.font.color.rgb = C_BLACK; r.font.name = FONT
            iy += Inches(0.92)


# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    builders = [
        (s01_cover,              "Cover"),
        (s02_current_state,      "현황 & 과제"),
        (s03_before_after,       "Before vs After"),
        (s04_writing_process,    "TC 작성 프로세스"),
        (s05_management_structure, "TC 관리 구조"),
        (s06_automation_process, "TC 자동화 개발 프로세스"),
        (s07_execution,          "TC 실행 & 결과 관리"),
        (s08_dashboard,          "대시보드 기능"),
        (s09_requirements,       "시스템 요구사항"),
        (s10_tech_stack,         "기술 스택"),
        (s11_data_structure,     "TC 데이터 구조"),
        (s12_change_management,  "변경 관리"),
        (s13_conclusion,         "결론"),
        (s14_next_steps,         "Next Steps"),
    ]

    for fn, name in builders:
        fn(prs)
        print(f"  ✓ {name}")

    today = datetime.date.today().strftime("%Y%m%d")
    out_dir = "/Users/myungjookim/_claude26/pptx"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"SPCY_TC_Proposal_Presentation_{today}_v2.pptx")
    prs.save(out_path)
    print(f"\n✅  저장 완료: {out_path}")
    return out_path


if __name__ == "__main__":
    build()
